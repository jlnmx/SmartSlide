from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

def create_business_template(output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Business Presentation"
    subtitle.text = "Subtitle for Business Template"
    prs.save(output_path)

def create_education_template(output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Education Presentation"
    subtitle.text = "Subtitle for Education Template"
    prs.save(output_path)

def create_creative_template(output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Creative Presentation"
    subtitle.text = "Subtitle for Creative Template"

    # Add a triangle shape for decoration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.5), Inches(4), Inches(2), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 102, 102)  # Red
    prs.save(output_path)

def create_modern_template(output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Modern Presentation"
    subtitle.text = "Subtitle for Modern Template"

    # Add a rectangle shape for a modern look
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 153, 76)  # Green
    prs.save(output_path)

def create_abstract_template(output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Abstract Presentation"
    subtitle.text = "Subtitle for Abstract Template"

    # Add a circle shape for an abstract look
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(102, 0, 204)  # Purple
    prs.save(output_path)

def create_minimal_template(output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Minimal Presentation"
    subtitle.text = "Subtitle for Minimal Template"

    # Add a simple line for a minimalistic look
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    line = slide.shapes.add_connector(
        connector_type=1,  # 1 corresponds to a straight line
        begin_x=Inches(1),
        begin_y=Inches(1),
        end_x=Inches(5),
        end_y=Inches(1)
    )
    line.line.color.rgb = RGBColor(0, 0, 0)  # Black
    prs.save(output_path)
    
# Generate the templates
create_business_template("templates/business_template.pptx")
create_education_template("templates/education_template.pptx")
create_creative_template("templates/creative_template.pptx")
create_modern_template("templates/modern_template.pptx")
create_abstract_template("templates/abstract_template.pptx")
create_minimal_template("templates/minimal_template.pptx")