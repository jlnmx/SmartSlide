import os
import re
import json
import requests
import uuid
import base64
import pandas as pd
import traceback
import firebase_admin
import os
import re
import json
import requests
import uuid
import base64
import pandas as pd
import traceback
import firebase_admin
from firebase_admin import credentials, firestore, auth
from werkzeug.security import generate_password_hash, check_password_hash
from datetime import datetime, timedelta, timezone
from flask import Blueprint, request, jsonify, current_app, send_from_directory, send_file # Added send_file
from dotenv import load_dotenv
from flask_cors import CORS
from werkzeug.utils import secure_filename
from io import BytesIO
from pptx import Presentation as PptxPresentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE # Ensure MSO_AUTO_SIZE is imported
from pdf2image import convert_from_path
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from google.oauth2 import service_account
from googleapiclient.discovery import build
from docx import Document # Added for Word export
from bs4 import BeautifulSoup
try:
    from youtube_transcript_api import YouTubeTranscriptApi
except ImportError:
    YouTubeTranscriptApi = None
from collections import Counter # Added import for Counter
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import random
import string
import secrets

# --- FIREBASE ADMIN INIT ---
cred = credentials.Certificate("firebase_key.json")
if not firebase_admin._apps:
    firebase_admin.initialize_app(cred)
firestore_db = firestore.client()

def update_analytics_on_slide(user_id, topic=None):
    """Update analytics for a user when a slide is generated."""
    analytics_ref = firestore_db.collection('analytics').document(user_id)
    analytics_doc = analytics_ref.get()
    if analytics_doc.exists:
        analytics_ref.update({
            'slides_created': firestore.Increment(1),
            'last_topic': topic,
            'last_generated_at': firestore.SERVER_TIMESTAMP
        })
    else:
        analytics_ref.set({
            'user_id': user_id,
            'slides_created': 1,
            'last_topic': topic,
            'last_generated_at': firestore.SERVER_TIMESTAMP
        })

def hex_to_rgb(hex_color, default_color="bfdbfe"):
    """Convert hex color to RGB tuple. Expands 3-digit hex, defaults to the provided default_color if invalid."""
    if not hex_color:
        current_app.logger.warning(f"Empty hex color, using default {default_color}")
        hex_color = default_color
    if hex_color.startswith('#'):
        hex_color = hex_color[1:]
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    if len(hex_color) != 6:
        current_app.logger.warning(f"Invalid hex color length: '{hex_color}', using default {default_color}")
        hex_color = default_color
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except ValueError as e:
        current_app.logger.warning(f"Invalid hex color: '{hex_color}', using default {default_color}. Error: {e}")
        return tuple(int(default_color[i:i+2], 16) for i in (0, 2, 4))

def calculate_position(shape_def, position_type, slide_width_inches, slide_height_inches):
    """Calculate position (left/top) based on shape definition"""
    # Check for direct inch values first
    inch_key = f"{position_type}_in"
    if inch_key in shape_def:
        return Inches(shape_def[inch_key])
        
    # Check for ratio-based positioning
    ratio_key = f"{position_type}_ratio"
    if ratio_key in shape_def:
        ratio = shape_def[ratio_key]
        if position_type == "left":
            return Inches(ratio * slide_width_inches)
        else:  # top
            return Inches(ratio * slide_height_inches)
            
    # Default to 0
    return Inches(0)

def calculate_dimension(shape_def, dimension_type, slide_width_inches, slide_height_inches):
    """Calculate dimension (width/height) based on shape definition"""
    # Check for direct inch values first
    inch_key = f"{dimension_type}_in"
    if inch_key in shape_def:
        return Inches(shape_def[inch_key])
        
    # Check for ratio-based sizing
    ratio_key = f"{dimension_type}_ratio"
    if ratio_key in shape_def:
        ratio = shape_def[ratio_key]
        if dimension_type == "width":
            return Inches(ratio * slide_width_inches)
        else:  # height
            return Inches(ratio * slide_height_inches)
            
    # Check for slide-height-based sizing
    height_ratio_key = f"{dimension_type}_as_ratio_of_slide_height"
    if height_ratio_key in shape_def:
        ratio = shape_def[height_ratio_key]
        return Inches(ratio * slide_height_inches)
        
    # Default to 1 inch
    return Inches(1)

def draw_template_shape(slide, shape_def, slide_width_inches, slide_height_inches):
    """
    Draw a shape on the slide based on shape definition
    
    Args:
        slide: PowerPoint slide object
        shape_def: Shape definition from template
        slide_width_inches: Slide width in inches
        slide_height_inches: Slide height in inches
    
    Returns:
        bool: True if shape was drawn successfully, False otherwise
    """
    try:
        if not shape_def:
            current_app.logger.warning("Empty shape definition provided")
            return False
            
        shape_type = shape_def.get("type", "rectangle")
        current_app.logger.debug(f"Drawing template shape: type={shape_type}, def={shape_def}")
        
        # Calculate position and size with validation
        left = calculate_position(shape_def, "left", slide_width_inches, slide_height_inches)
        top = calculate_position(shape_def, "top", slide_width_inches, slide_height_inches)
        width = calculate_dimension(shape_def, "width", slide_width_inches, slide_height_inches)
        height = calculate_dimension(shape_def, "height", slide_width_inches, slide_height_inches)
        
        # Validate dimensions
        if width <= 0 or height <= 0:
            current_app.logger.warning(f"Invalid shape dimensions: width={width}, height={height}")
            return False
        
        from pptx.enum.shapes import MSO_SHAPE
        
        # Use rounded rectangle for content box if specified
        is_content_box = shape_def.get("comment", "") == "content_box"
        if shape_type == "rectangle" and is_content_box:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            # Set rounding if radius is specified (python-pptx uses adjustment values 0-1)
            radius = shape_def.get("radius")
            if radius is not None:
                try:
                    # Adjustment 0 is corner rounding, 0.0 (square) to 1.0 (fully round)
                    shape.adjustments[0] = float(radius)
                except Exception as radius_error:
                    current_app.logger.warning(f"Failed to set shape radius: {radius_error}")
        elif shape_type == "rectangle":
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        elif shape_type == "oval":
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        else:
            current_app.logger.warning(f"Unsupported shape type: {shape_type}")
            return False
            
        # Apply fill
        fill_style = shape_def.get("fill_style", "solid")
        try:
            if fill_style == "none":
                shape.fill.background()
            else:
                fill_color_hex = shape_def.get("fill_color_hex", "FFFFFF")
                fill_transparency = shape_def.get("fill_transparency", 0.0)
                
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color_hex))
                if fill_transparency > 0:
                    shape.fill.transparency = fill_transparency
        except Exception as fill_error:
            current_app.logger.warning(f"Failed to apply shape fill: {fill_error}")
                
        # Apply line/border
        line_style = shape_def.get("line_style", "solid")
        try:
            if line_style == "none":
                shape.line.fill.background()
            else:
                line_color_hex = shape_def.get("line_color_hex", "000000")
                line_width_pt = shape_def.get("line_width_pt", 1)
                line_transparency = shape_def.get("line_transparency", 0.0)
                
                shape.line.color.rgb = RGBColor(*hex_to_rgb(line_color_hex))
                shape.line.width = Pt(line_width_pt)
                if line_transparency > 0:
                    shape.line.transparency = line_transparency
        except Exception as line_error:
            current_app.logger.warning(f"Failed to apply shape line: {line_error}")
                
        # Apply rotation if specified
        rotation = shape_def.get("rotation", 0)
        if rotation != 0:
            try:
                shape.rotation = rotation
            except Exception as rotation_error:
                current_app.logger.warning(f"Failed to apply shape rotation: {rotation_error}")
        
        current_app.logger.debug(f"Successfully drew template shape: {shape_type}")
        return True
            
    except Exception as e:
        current_app.logger.error(f"Error drawing template shape: {e}")
        return False


load_dotenv()
GOOGLE_CREDENTIALS_FILE = os.path.join(os.path.dirname(__file__), '..', 'credentials.json')
GOOGLE_CREDENTIALS_FILE = os.path.abspath(GOOGLE_CREDENTIALS_FILE)
SCOPES = ["https://www.googleapis.com/auth/presentations", "https://www.googleapis.com/auth/drive.file"]
ALLOWED_EXTENSIONS = {'pdf', 'xlsx', 'xls', 'csv', 'docx'}
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
if not GROQ_API_KEY:
    raise RuntimeError("GROQ_API_KEY environment variable is not set. Please set it in your environment or .env file.")

# Email configuration for SMTP
SMTP_SERVER = os.environ.get("SMTP_SERVER", "smtp.gmail.com")
SMTP_PORT = int(os.environ.get("SMTP_PORT", "587"))
EMAIL_ADDRESS = os.environ.get("EMAIL_ADDRESS")
EMAIL_PASSWORD = os.environ.get("EMAIL_PASSWORD")

def send_email_verification(email, verification_code):
    """Send email verification code using Firebase Auth and SMTP fallback"""
    try:
        # Option 1: Use Firebase Auth's built-in email verification
        # Note: This requires Firebase project configuration for email templates
        
        # Option 2: Send via SMTP (Gmail)
        if not EMAIL_ADDRESS or not EMAIL_PASSWORD:
            current_app.logger.info(f"Email service not configured, logging code for {email}: {verification_code}")
            return True, "Email service configured (development mode)"
        
        # Create message
        msg = MIMEMultipart()
        msg['From'] = EMAIL_ADDRESS
        msg['To'] = email
        msg['Subject'] = "SmartSlide Password Reset Verification Code"
        
        # Email body with better formatting
        body = f"""
        <html>
        <body>
            <div style="max-width: 600px; margin: 0 auto; font-family: Arial, sans-serif;">
                <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 20px; text-align: center;">
                    <h1 style="color: white; margin: 0;">SmartSlide</h1>
                    <p style="color: white; margin: 5px 0;">Password Reset Request</p>
                </div>
                
                <div style="padding: 30px; background: #f9f9f9;">
                    <h2 style="color: #333;">Password Reset Verification</h2>
                    <p style="color: #666; line-height: 1.6;">
                        You requested a password reset for your SmartSlide account. Use the verification code below to proceed:
                    </p>
                    
                    <div style="background: white; padding: 20px; border-radius: 8px; text-align: center; margin: 20px 0;">
                        <h3 style="color: #333; margin: 0;">Verification Code</h3>
                        <div style="font-size: 32px; font-weight: bold; color: #667eea; letter-spacing: 8px; margin: 15px 0;">
                            {verification_code}
                        </div>
                    </div>
                    
                    <p style="color: #666; line-height: 1.6;">
                        This code will expire in <strong>10 minutes</strong>.
                    </p>
                    
                    <p style="color: #666; line-height: 1.6;">
                        If you didn't request this password reset, please ignore this email. Your account remains secure.
                    </p>
                </div>
                
                <div style="background: #333; padding: 15px; text-align: center;">
                    <p style="color: #ccc; margin: 0; font-size: 12px;">
                        © 2025 SmartSlide. All rights reserved.
                    </p>
                </div>
            </div>
        </body>
        </html>
        """
        
        msg.attach(MIMEText(body, 'html'))
        
        # Send email
        server = smtplib.SMTP(SMTP_SERVER, SMTP_PORT)
        server.starttls()
        server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
        text = msg.as_string()
        server.sendmail(EMAIL_ADDRESS, email, text)
        server.quit()
        
        current_app.logger.info(f"Email sent successfully to {email}")
        return True, "Email sent successfully"
        
    except Exception as e:
        current_app.logger.error(f"Failed to send email to {email}: {e}")
        # Fallback: log the code for development
        current_app.logger.info(f"Email service failed, password reset code for {email}: {verification_code}")
        return True, "Email sent (development mode)"

def send_sms_verification(phone_number, verification_code):
    """Send SMS verification code using Firebase Auth"""
    try:
        # Format phone number for international format
        formatted_phone = phone_number
        if not phone_number.startswith('+'):
            # Assume Philippines if no country code
            if phone_number.startswith('09'):
                formatted_phone = '+63' + phone_number[1:]
            elif phone_number.startswith('9'):
                formatted_phone = '+63' + phone_number
            else:
                formatted_phone = '+63' + phone_number
        
        # For development: Log the SMS code
        current_app.logger.info(f"Password reset SMS code for {formatted_phone}: {verification_code}")
        
        # Option 1: Use Firebase Auth Phone Authentication
        # Note: This requires Firebase project configuration for SMS
        # You can enable this in Firebase Console > Authentication > Sign-in method > Phone
        
        # Option 2: Use a free SMS service or development logging
        # For now, we'll log the code for development
        
        return True, f"SMS sent to {formatted_phone} (development mode)"
        
    except Exception as e:
        current_app.logger.error(f"Failed to send SMS to {phone_number}: {e}")
        current_app.logger.info(f"SMS service failed, password reset code for {phone_number}: {verification_code}")
        return True, "SMS sent (development mode)"

main = Blueprint("main", __name__)

# Add proper CORS headers for all responses
@main.after_request
def after_request(response):
    origin = request.headers.get('Origin')
    # Allow localhost for development and all Vercel deployments
    allowed_origins = [
        'http://localhost:3000',
        'https://smartslide.vercel.app',
        'https://smartslide-git-main-your-username.vercel.app',  # Git branch deployments
        'https://smartslide-git-master-your-username.vercel.app'  # Master branch deployments
    ]
    
    # Also allow any vercel.app subdomain for your project
    if origin and (origin in allowed_origins or (origin.endswith('.vercel.app') and 'smartslide' in origin)):
        response.headers.add('Access-Control-Allow-Origin', origin)
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization,X-Requested-With,Accept,Origin')
        response.headers.add('Access-Control-Allow-Methods', 'GET,POST,PUT,DELETE,OPTIONS,PATCH')
        response.headers.add('Access-Control-Allow-Credentials', 'true')
        response.headers.add('Access-Control-Max-Age', '86400')
        response.headers.add('Vary', 'Origin')
    return response

# Handle preflight requests globally
@main.before_request
def handle_preflight():
    if request.method == "OPTIONS":
        origin = request.headers.get('Origin')
        # Allow localhost for development and all Vercel deployments
        allowed_origins = [
            'http://localhost:3000',
            'https://smartslide.vercel.app',
            'https://smartslide-git-main-your-username.vercel.app',
            'https://smartslide-git-master-your-username.vercel.app'
        ]
        
        # Also allow any vercel.app subdomain for your project
        if origin and (origin in allowed_origins or (origin.endswith('.vercel.app') and 'smartslide' in origin)):
            response = jsonify({'status': 'ok'})
            response.headers.add('Access-Control-Allow-Origin', origin)
            response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization,X-Requested-With,Accept,Origin')
            response.headers.add('Access-Control-Allow-Methods', 'GET,POST,PUT,DELETE,OPTIONS,PATCH')
            response.headers.add('Access-Control-Allow-Credentials', 'true')
            response.headers.add('Access-Control-Max-Age', '86400')
            response.headers.add('Vary', 'Origin')
            return response

# --- SIMPLE HEALTH CHECK FOR CORS TESTING ---
@main.route('/health', methods=['GET', 'OPTIONS'])
def health_check():
    return jsonify({'status': 'ok', 'message': 'Backend is running'}), 200

# --- FIREBASE USER REGISTRATION ---
@main.route('/register', methods=['POST', 'OPTIONS'])
def register():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.get_json()
    email = data.get('email')
    password = data.get('password')
    full_name = data.get('fullName')
    birthday = data.get('birthday')
    contact_number = data.get('contactNumber')
    
    # Validate required fields
    if not email or not password or not full_name or not birthday or not contact_number:
        return jsonify({'error': 'All fields are required (email, password, fullName, birthday, contactNumber)'}), 400
    
    users_ref = firestore_db.collection('users')
    existing = users_ref.where('email', '==', email).get()
    if existing:
        return jsonify({'error': 'Email already registered'}), 409
    
    user_doc = users_ref.document()
    user_doc.set({
        'email': email,
        'password_hash': generate_password_hash(password),
        'full_name': full_name,
        'birthday': birthday,
        'contact_number': contact_number,
        'user_type': None,  # Will be set in the user type selection step
        'registration_completed': False,  # Will be set to True after user type selection
        'created_at': firestore.SERVER_TIMESTAMP
    })
    return jsonify({'message': 'User registered', 'user_id': user_doc.id}), 201

# --- FIREBASE LOGIN ---
@main.route('/login', methods=['POST', 'OPTIONS'])
def login():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    data = request.get_json()
    email = data.get('email')
    password = data.get('password')
    users_ref = firestore_db.collection('users')
    user_query = users_ref.where('email', '==', email).get()    
    if not user_query:
        return jsonify({'error': 'Invalid email or password'}), 401
    user = user_query[0]
    if not check_password_hash(user.get('password_hash'), password):
        return jsonify({'error': 'Invalid email or password'}), 401
    
    # Return user data including registration status
    user_data = user.to_dict()
    return jsonify({
        'message': 'Login successful', 
        'user': {
            'id': user.id, 
            'email': user_data.get('email'),
            'full_name': user_data.get('full_name'),
            'user_type': user_data.get('user_type'),
            'registration_completed': user_data.get('registration_completed', False)
        }
    }), 200

# --- USER TYPE SELECTION ---
@main.route('/select-user-type', methods=['POST', 'OPTIONS'])
def select_user_type():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.get_json()
    user_id = data.get('user_id')
    user_type = data.get('user_type')
    
    # Validate inputs
    if not user_id or not user_type:
        return jsonify({'error': 'User ID and user type are required'}), 400
    
    if user_type not in ['student', 'professional', 'personal']:
        return jsonify({'error': 'Invalid user type. Must be student, professional, or personal'}), 400
    
    try:
        # Update user document with user type and mark registration as completed
        users_ref = firestore_db.collection('users')
        user_doc_ref = users_ref.document(user_id)
        
        # Check if user exists
        user_doc = user_doc_ref.get()
        if not user_doc.exists:
            return jsonify({'error': 'User not found'}), 404
        
        # Update user type and registration status
        user_doc_ref.update({
            'user_type': user_type,
            'registration_completed': True,
            'updated_at': firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            'message': 'User type selected successfully',
            'user_type': user_type
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Error updating user type: {e}")
        return jsonify({'error': 'Failed to update user type'}), 500

# --- FIREBASE USER PROFILE ---
@main.route('/user/<user_id>', methods=['GET', 'PUT', 'OPTIONS'])
def manage_user_profile(user_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    users_ref = firestore_db.collection('users')
    user_doc = users_ref.document(user_id).get()
    if not user_doc.exists:
        return jsonify({'error': 'User not found'}), 404
    if request.method == 'GET':
        user = user_doc.to_dict()
        user_data = {
            'id': user_id,
            'email': user.get('email'),
            'full_name': user.get('full_name'),
            'birthday': user.get('birthday'),
            'contact_number': user.get('contact_number'),
            'user_type': user.get('user_type'),
            'registration_completed': user.get('registration_completed', False)
        }
        return jsonify(user_data), 200
    if request.method == 'PUT':
        data = request.get_json()
        update_data = {}
        if 'email' in data:
            update_data['email'] = data['email']
        if 'full_name' in data:
            update_data['full_name'] = data['full_name']
        if 'birthday' in data:
            update_data['birthday'] = data['birthday']
        if 'contact_number' in data:
            update_data['contact_number'] = data['contact_number']
        if 'password' in data and data['password']:
            update_data['password_hash'] = generate_password_hash(data['password'])
        if update_data:
            update_data['updated_at'] = firestore.SERVER_TIMESTAMP
            users_ref.document(user_id).update(update_data)
        user_doc = users_ref.document(user_id).get()
        user = user_doc.to_dict()
        user_data = {
            'id': user_id,
            'email': user.get('email'),
            'full_name': user.get('full_name'),
            'birthday': user.get('birthday'),
            'contact_number': user.get('contact_number'),
            'user_type': user.get('user_type'),
            'registration_completed': user.get('registration_completed', False)
        }
        return jsonify({'message': 'User profile updated successfully', 'user': user_data}), 200
    return jsonify({'error': 'Invalid request method'}), 405

# --- FIREBASE SAVE PRESENTATION ---
@main.route('/save-presentation', methods=['POST', 'OPTIONS'])
def save_presentation():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    data = request.get_json()
    if not data:
        return jsonify({'error': 'No data provided'}), 400
    user_id = data.get('user_id')
    title = data.get('title', 'Untitled Presentation')
    slides = data.get('slides')
    template = data.get('template', 'default')
    presentation_type = data.get('presentation_type', 'Default')
    if not user_id or not slides:
        return jsonify({'error': 'Missing required fields (user_id, slides)'}), 400
    try:
        now = datetime.utcnow()
        doc = firestore_db.collection('presentations').document()
        doc.set({
            'user_id': user_id,
            'title': title,
            'slides': slides,
            'template': template,
            'presentation_type': presentation_type,
            'created_at': firestore.SERVER_TIMESTAMP,
            'updated_at': firestore.SERVER_TIMESTAMP
        })
        return jsonify({'message': 'Presentation saved successfully', 'presentationId': doc.id}), 201
    except Exception as e:
        current_app.logger.error(f"Error saving presentation: {e}", exc_info=True)
        return jsonify({'error': 'Failed to save presentation'}), 500

# --- FIREBASE GET PRESENTATIONS FOR USER ---
@main.route('/presentations/<user_id>', methods=['GET', 'OPTIONS'])
def get_presentations(user_id):
    """Return the 10 most recent presentations for a user, with ISO-formatted timestamps."""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
        
    try:
        current_app.logger.info(f"Fetching presentations for user_id: {user_id}")
        presentations_ref = firestore_db.collection('presentations')
        query = presentations_ref.where('user_id', '==', user_id)
        docs = query.stream()
        presentations = []
        for doc in docs:
            data = doc.to_dict()
            data['id'] = doc.id
            for ts_field in ['created_at', 'updated_at']:
                if ts_field in data and data[ts_field] is not None:
                    try:
                        data[ts_field] = data[ts_field].isoformat()
                    except Exception:
                        data[ts_field] = str(data[ts_field])
            presentations.append(data)
        presentations.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        presentations = presentations[:10]
        current_app.logger.info(f"Returning {len(presentations)} presentations for user_id: {user_id}")
        return jsonify({'presentations': presentations}), 200
    except Exception as e:
        current_app.logger.error(f"Error fetching presentations for user_id {user_id}: {e}")
        return jsonify({'presentations': [], 'error': str(e)}), 200

# --- FIREBASE SAVE QUIZ ---
@main.route('/save-quiz', methods=['POST', 'OPTIONS'])
def save_quiz():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    data = request.get_json()
    user_id = data.get('user_id')
    name = data.get('name')
    content = data.get('content')
    if not user_id or not name or not content:
        return jsonify({'error': 'Missing required fields'}), 400
    quizzes_ref = firestore_db.collection('saved_quizzes')
    doc = quizzes_ref.document()
    doc.set({
        'user_id': user_id,
        'name': name,
        'content': content,
        'created_at': firestore.SERVER_TIMESTAMP
    })
    return jsonify({'message': 'Quiz saved', 'quiz_id': doc.id}), 201

# --- FIREBASE SAVE SCRIPT ---
@main.route('/save-script', methods=['POST', 'OPTIONS'])
def save_script():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    data = request.get_json()
    user_id = data.get('user_id')
    name = data.get('name')
    content = data.get('content')
    if not user_id or not name or not content:
        return jsonify({'error': 'Missing required fields'}), 400
    scripts_ref = firestore_db.collection('saved_scripts')
    doc = scripts_ref.document()
    doc.set({
        'user_id': user_id,
        'name': name,
        'content': content,
        'created_at': firestore.SERVER_TIMESTAMP
    })
    return jsonify({'message': 'Script saved', 'script_id': doc.id}), 201

# --- FIREBASE GET SAVED QUIZZES & SCRIPTS FOR USER ---
@main.route('/saved-items/<user_id>', methods=['GET', 'OPTIONS'])
def get_saved_items(user_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    try:
        quizzes_ref = firestore_db.collection('saved_quizzes').where('user_id', '==', user_id)
        scripts_ref = firestore_db.collection('saved_scripts').where('user_id', '==', user_id)
        quizzes = [dict(q.to_dict(), id=q.id) for q in quizzes_ref.stream()]
        scripts = [dict(s.to_dict(), id=s.id) for s in scripts_ref.stream()]
        quizzes.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        scripts.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        return jsonify({'quizzes': quizzes, 'scripts': scripts}), 200
    except Exception as e:
        current_app.logger.error(f"Error fetching saved items for user_id {user_id}: {e}")
        return jsonify({'quizzes': [], 'scripts': [], 'error': str(e)}), 500

# --- FIREBASE DELETE QUIZ ---
@main.route('/saved-quiz/<quiz_id>', methods=['DELETE', 'OPTIONS'])
def delete_saved_quiz(quiz_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    firestore_db.collection('saved_quizzes').document(quiz_id).delete()
    return jsonify({'message': 'Quiz deleted successfully'}), 200

# --- FIREBASE DELETE SCRIPT ---
@main.route('/saved-script/<script_id>', methods=['DELETE', 'OPTIONS'])
def delete_saved_script(script_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    firestore_db.collection('saved_scripts').document(script_id).delete()
    return jsonify({'message': 'Script deleted successfully'}), 200

@main.route('/generate-quiz', methods=['POST', 'OPTIONS'])
def generate_quiz_route():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.json
    slides_content = data.get("slides")
    if not slides_content:
        return jsonify({"error": "No slides content provided"}), 400    # Combine slide content into a single text for the prompt
    full_text_content = ""
    for slide in slides_content:
        # Handle both editor format (textboxes) and simple format (title/content)
        title = slide.get("title") or ""
        
        # Check if this is from the slide editor (has textboxes)
        if slide.get("textboxes") and isinstance(slide["textboxes"], list):
            title_box = next((tb for tb in slide["textboxes"] if tb.get("type") == "title"), None)
            body_box = next((tb for tb in slide["textboxes"] if tb.get("type") == "body"), None)
            
            if title_box:
                title = title_box.get("text", title)
            
            content = ""
            if body_box:
                content = body_box.get("text", "")
            
            # Also collect any other textboxes that might contain content
            other_textboxes = [tb for tb in slide["textboxes"] if tb.get("type") not in ["title", "body"] and tb.get("text")]
            for tb in other_textboxes:
                if content:
                    content += "\n"
                content += tb.get("text", "")
        else:
            # Handle simple format
            content = slide.get("content", "")
            if isinstance(content, list):
                content = "\n".join(content)
            elif not isinstance(content, str):
                content = str(content)
        
        if title:
            full_text_content += f"Slide Title: {title}\n"        
        if content:
            full_text_content += f"{content}\n\n"

    if not full_text_content.strip():
        current_app.logger.error(f"No content extracted from slides. Slides data: {slides_content}")
        return jsonify({"error": "No content found in slides"}), 400

    current_app.logger.info(f"Extracted content for quiz generation: {full_text_content[:500]}...")  # Log first 500 chars

    # Get language and number of questions from the request
    language = data.get("language", "English")  # Default to English
    num_questions = int(data.get("numQuestions", 5))  # Default to 5 questions

    quiz_prompt = f"""
    Based on the following presentation content, generate a quiz with exactly {num_questions} questions.
    The quiz should include a mix of identification and multiple-choice questions.
    For multiple-choice questions, provide 4 choices.
    Provide a clear answer for each question.
    The language for the quiz must be {language}.
    Format the output as a JSON array, where each object has "question", "choices" (an array of 4 strings for multiple-choice, or an empty array/null for identification questions), and "answer" (a string).

    Presentation Content:
    ---
    {full_text_content}
    ---

    Generate the quiz now in JSON format:
    """
    try:
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "llama3-8b-8192",
            "messages": [
                {"role": "system", "content": "You are an assistant that generates quizzes in JSON format based on provided text."},
                {"role": "user", "content": quiz_prompt}
            ],
            "max_tokens": 2048,
            "temperature": 0.5
        }
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)
        response.raise_for_status()

        result = response.json()
        model_output = result["choices"][0]["message"]["content"]
        
        # Extract JSON array from the response
        match = re.search(r'\[\s*{.*}\s*\]', model_output, re.DOTALL)
        if not match:
            current_app.logger.error(f"Failed to parse quiz JSON from LLM output: {model_output}")
            return jsonify({"error": "Failed to generate valid quiz format"}), 500
        
        quiz_data = json.loads(match.group(0))
        return jsonify({"quiz": quiz_data})

    except requests.exceptions.HTTPError as http_err:
        error_details = "N/A"
        if http_err.response is not None:
            try:
                error_details = http_err.response.json()
            except ValueError:
                error_details = http_err.response.text
        current_app.logger.error(f"HTTP error during quiz generation: {http_err}. Details: {error_details}")
        return jsonify({"error": "Failed to communicate with AI service"}), 500    
    except json.JSONDecodeError as e:
        output_for_log = locals().get('model_output', 'N/A')
        current_app.logger.error(f"JSON Decode Error in quiz generation: {e}. LLM Output: {output_for_log}")
        return jsonify({"error": "Failed to parse quiz response"}), 500
    except Exception as e:
        current_app.logger.error(f"Error during quiz generation: {e}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error"}), 500

@main.route('/generate-script', methods=['POST', 'OPTIONS'])
def generate_script_route():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
        
    data = request.json
    slides_content = data.get("slides")
    if not slides_content:
        return jsonify({"error": "No slides content provided"}), 400

    # Combine slide content into a single text for the prompt
    full_text_content = ""
    for slide in slides_content:
        title = slide.get("title") or None
        if slide.get("textboxes") and isinstance(slide["textboxes"], list):
            title_box = next((tb for tb in slide["textboxes"] if tb.get("type") == "title"), None)
            body_box = next((tb for tb in slide["textboxes"] if tb.get("type") == "body"), None)
            if title_box:
                title = title_box.get("text", title or "")
            content = body_box.get("text", "") if body_box else ""
        else:
            content = slide.get("content", "")
            if isinstance(content, list):
                content = "\n".join(content)
            elif not isinstance(content, str):
                content = str(content)
        if title:
            full_text_content += f"Slide Title: {title}\n"
        if content:
            full_text_content += f"{content}\n\n"

    if not full_text_content.strip():
        return jsonify({"error": "No content found in slides"}), 400

    script_prompt = f"""
    Based on the following presentation content, generate a detailed speaker script.
    The script should elaborate on the key points of each slide, provide transitions, and suggest where to pause or emphasize.
    The output should be a single string of text.

    Presentation Content:
    ---
    {full_text_content}
    ---

    Generate the speaker script now:
    """
    try:
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "llama3-8b-8192",
            "messages": [
                {"role": "system", "content": "You are an assistant that generates speaker scripts based on provided presentation content."},
                {"role": "user", "content": script_prompt}
            ],
            "max_tokens": 3000,
            "temperature": 0.6
        }
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)
        response.raise_for_status()

        result = response.json()
        script_data = result["choices"][0]["message"]["content"]
        return jsonify({"script": script_data})

    except requests.exceptions.HTTPError as http_err:
        error_details = "N/A"
        if http_err.response is not None:
            try:
                error_details = http_err.response.json()
            except ValueError:
                error_details = http_err.response.text
        current_app.logger.error(f"HTTP error during script generation: {http_err}. Details: {error_details}")
        return jsonify({"error": "Failed to communicate with AI service"}), 500
    except Exception as e:
        current_app.logger.error(f"Error during script generation: {e}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error"}), 500

# # --- FIREBASE GET RECENT PRESENTATIONS FOR USER ---
# @main.route('/presentations/<user_id>', methods=['GET'])
# def get_recent_presentations(user_id):
#     """Return the most recent presentations for a user (limit 10)."""
#     presentations_ref = firestore_db.collection('presentations')
#     query = presentations_ref.where('user_id', '==', user_id).order_by('created_at', direction=firestore.Query.DESCENDING).limit(10)
#     docs = query.stream()
#     presentations = []
#     for doc in docs:
#         data = doc.to_dict()
#         data['id'] = doc.id
#         presentations.append(data)
#     return jsonify({'presentations': presentations})


# --- FIREBASE GET/DELETE SINGLE PRESENTATION ---
@main.route('/presentation/<presentation_id>', methods=['GET', 'DELETE', 'OPTIONS'])
def manage_presentation(presentation_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    presentation_ref = firestore_db.collection('presentations').document(str(presentation_id))

    if request.method == 'GET':
        try:
            presentation_doc = presentation_ref.get()
            if not presentation_doc.exists:
                return jsonify({'error': 'Presentation not found'}), 404
            pres_data = presentation_doc.to_dict()
            slides_data = pres_data.get('slides', [])
            return jsonify({
                'id': presentation_doc.id,
                'title': pres_data.get('title'),
                'slides': slides_data,
                'template': pres_data.get('template'),
                'presentation_type': pres_data.get('presentation_type'),
                'created_at': pres_data.get('created_at')
            })
        except Exception as e:
            current_app.logger.error(f"Error retrieving presentation {presentation_id}: {e}")
            return jsonify({'error': 'Failed to retrieve presentation'}), 500

    elif request.method == 'DELETE':
        try:
            presentation_ref.delete()
            return jsonify({'message': 'Presentation deleted successfully'}), 200
        except Exception as e:
            current_app.logger.error(f"Error deleting presentation {presentation_id}: {e}")
            return jsonify({'error': 'Failed to delete presentation'}), 500

# Route to save or update slide editor state
@main.route('/api/save-slides-state', methods=['POST', 'OPTIONS'])
def save_slides_state():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "Invalid JSON payload"}), 400

        user_id = data.get("userId") or data.get("user_id")
        slides = data.get("slides")
        template_id = data.get("templateId") or data.get("template")
        presentation_type = data.get("presentationType", "Default")
        presentation_id = data.get("presentationId") or data.get("presentation_id")

        if not user_id or not slides or not template_id:
            return jsonify({"error": "Missing required fields (user_id, slides, template_id)"}), 400

        slides_json_str = json.dumps(slides)
        now = datetime.utcnow()        
        
        if presentation_id:
            # Update existing presentation
            presentation_ref = firestore_db.collection('presentations').document(str(presentation_id))
            presentation_doc = presentation_ref.get()
            if not presentation_doc.exists:
                return jsonify({"error": "Presentation not found"}), 404
            
            # Update the title from the slides if it has changed
            if slides and isinstance(slides, list) and len(slides) > 0:
                first_slide = slides[0]
                new_title = None
                
                # Check if it's editor format (with textboxes)
                if isinstance(first_slide, dict) and 'textboxes' in first_slide:
                    title_textbox = next((tb for tb in first_slide['textboxes'] if tb.get('type') == 'title'), None)
                    if title_textbox and title_textbox.get('text'):
                        new_title = title_textbox['text']
                # Check if it's simple format (with title field)
                elif isinstance(first_slide, dict) and 'title' in first_slide:
                    new_title = first_slide['title']
                
                if new_title and new_title.strip() and new_title != "Title":
                    presentation_ref.update({'title': new_title.strip()})
            
            presentation_ref.update({
                'slides': slides,
                'template': template_id,
                'presentation_type': presentation_type,
                'updated_at': now
            })
            return jsonify({"message": "Presentation updated successfully", "presentationId": presentation_id}), 200
        else:
            # Create new presentation
            title = "Untitled Presentation"
            if slides and isinstance(slides, list) and len(slides) > 0:
                first_slide = slides[0]
                
                # Check if it's editor format (with textboxes)
                if isinstance(first_slide, dict) and 'textboxes' in first_slide:
                    title_textbox = next((tb for tb in first_slide['textboxes'] if tb.get('type') == 'title'), None)
                    if title_textbox and title_textbox.get('text') and title_textbox['text'].strip() != "Title":
                        title = title_textbox['text'].strip()
                # Check if it's simple format (with title field)
                elif isinstance(first_slide, dict) and 'title' in first_slide and first_slide['title'].strip():
                    title = first_slide['title'].strip()
            
            doc = firestore_db.collection('presentations').document()
            doc.set({
                'user_id': user_id,
                'title': title,
                'slides': slides,
                'template': template_id,
                'presentation_type': presentation_type,
                'created_at': now,
                'updated_at': now
            })
            return jsonify({"message": "Presentation created successfully", "presentationId": doc.id}), 201
    except Exception as e:
        current_app.logger.error(f"Error in /api/save-slides-state: {e}", exc_info=True)
        return jsonify({"error": "An error occurred while saving the presentation."}), 500

# --- FIREBASE GENERATE SLIDES ENDPOINT (REMOVE SQLAlchemy) ---
@main.route("/generate-slides", methods=["POST", "OPTIONS"])
def generate_slides():
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'ok'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
        response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
        return response
        
    data = request.json
    if not data:
        return jsonify({"error": "No data provided"}), 400

    prompt_topic = data.get("prompt")
    language = data.get("language", "English")
    user_id = data.get("user_id")  # Expect user_id from frontend
    template = data.get("template")
    try:
        num_slides = int(data.get("numSlides", 5))
        # Fix: Change validation to properly support up to 30 slides
        if num_slides <= 0 or num_slides > 30:
            return jsonify({"error": "Number of slides must be between 1 and 30."}), 400
    except (ValueError, TypeError):
        return jsonify({"error": "Invalid number of slides."}), 400
    if not prompt_topic:
        return jsonify({"error": "Prompt topic is required."}), 400
    try:
        import requests, json, re, traceback
        
        # Calculate content slides (excluding title, conclusion, and references)
        content_slides = max(1, num_slides - 3)  # At least 1 content slide
        
        generation_prompt = f"""
        Generate a professional, well-structured presentation about \"{prompt_topic}\".
        Requirements:
        - The presentation must have exactly {num_slides} slides.
        - Use {language} as the language.
        - Each slide should have:
          - A concise and engaging title.
          - Remove the unnecessary characters from the texts like the ** or __. 
          - Make the structure of the sentences or paragraph clean
          - Generate the best answers possible for the given topic and do not be frugal with the content.
          - Clear and concise content, formatted as bullet points or short paragraphs.
          - Use Markdown for formatting: **bold** for emphasis, *italic* for highlights, and __underline__ for key terms.
          - If applicable, include a relevant image description for each slide (e.g., "A diagram of the water cycle").
        
        - Organize the content logically:
          - Slide 1: Title slide with the topic and description.
          - Slide 2: Introduction (overview of the topic).
          - Slides 3-{num_slides-2}: Main content slides with detailed information, examples, case studies, analysis, etc.
          - Slide {num_slides-1}: Conclusion and Next Steps
          - Slide {num_slides}: References and Sources
        
        - For presentations with many slides ({num_slides} slides), ensure each content slide covers a specific aspect:
          - Background/History
          - Key Concepts/Definitions
          - Current State/Market Analysis
          - Trends and Developments
          - Challenges and Opportunities
          - Case Studies/Examples
          - Best Practices
          - Implementation Strategies
          - Risk Assessment
          - Future Outlook
          - Recommendations
          - Action Items
        
        - Ensure the content is professional, insightful, and suitable for a business or academic audience.
        - Each slide must have substantial content - no empty or placeholder slides.
        - Provide a JSON array where each object has:
          - 'title': string,
          - 'content': array of strings (use Markdown for formatting),
          - 'image_prompt': string (optional, for generating relevant images).

        Example structure for a {num_slides}-slide presentation:
        [
          {{
            "title": "{prompt_topic}",
            "content": [
              "A comprehensive overview of {prompt_topic}",
              "Understanding the key aspects and implications",
              "Strategic insights and analysis"
            ]
          }},
          {{
            "title": "Introduction and Overview",
            "content": [
              "Definition and scope of {prompt_topic}",
              "Why this topic is important in today's context",
              "Key questions we will address",
              "Expected outcomes from this presentation"
            ]
          }},
          {{
            "title": "Historical Background and Context",
            "content": [
              "Origins and evolution of {prompt_topic}",
              "Key milestones and developments",
              "Historical significance and impact",
              "Lessons learned from the past"
            ]
          }},
          {{
            "title": "Current Market Analysis",
            "content": [
              "Present state of {prompt_topic}",
              "Market size and growth trends",
              "Key players and stakeholders",
              "Current challenges and opportunities"
            ]
          }},
          {{
            "title": "Key Trends and Developments",
            "content": [
              "Emerging trends in {prompt_topic}",
              "Technological advancements",
              "Industry innovations",
              "Future growth projections"
            ]
          }},
          {{
            "title": "Challenges and Risk Factors",
            "content": [
              "Major obstacles and barriers",
              "Risk assessment and mitigation",
              "Common pitfalls to avoid",
              "Strategic considerations"
            ]
          }},
          {{
            "title": "Case Studies and Examples",
            "content": [
              "Real-world applications of {prompt_topic}",
              "Success stories and best practices",
              "Lessons from failures",
              "Industry benchmarks"
            ]
          }},
          {{
            "title": "Implementation Strategies",
            "content": [
              "Step-by-step approach to {prompt_topic}",
              "Resource requirements and planning",
              "Timeline and milestones",
              "Success metrics and KPIs"
            ]
          }},
          {{
            "title": "Conclusion and Next Steps",
            "content": [
              "Summary of key findings and insights",
              "Strategic recommendations",
              "Immediate action items",
              "Long-term considerations",
              "Follow-up activities and monitoring"
            ]
          }},
          {{
            "title": "References and Sources",
            "content": [
              "1. Academic journals and research papers",
              "2. Industry reports and white papers",
              "3. Expert interviews and surveys",
              "4. Government and regulatory publications",
              "5. Professional associations and standards"
            ]
          }}
        ]

        Generate exactly {num_slides} slides now with substantial content for each slide:
        """

        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        
        # Increase token limit for larger presentations
        max_tokens = 4096
        if num_slides > 15:
            max_tokens = 6144  # Increase token limit for larger presentations
        elif num_slides > 25:
            max_tokens = 8192  # Even more tokens for very large presentations
            
        payload = {
            "model": "llama3-8b-8192",  
            "messages": [
                {"role": "system", "content": "You are a helpful assistant that generates comprehensive slide content in JSON format. Always ensure the last two slides are Conclusion and References respectively."},
                {"role": "user", "content": generation_prompt}
            ],
            "max_tokens": max_tokens,
            "temperature": 0.3
        }
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)
        if response.status_code != 200:
            return jsonify({"error": "Failed to generate slides."}), 500

        result = response.json()
        model_output = result["choices"][0]["message"]["content"]

        # Extract JSON array from the response
        match = re.search(r'\[\s*{.*}\s*\]', model_output, re.DOTALL)
        if not match:
            return jsonify({"error": "Failed to parse slides output."}), 500
        slides_data = json.loads(match.group(0))

        # Validate and ensure we have the right number of slides with proper structure
        if len(slides_data) != num_slides:
            current_app.logger.warning(f"Generated {len(slides_data)} slides instead of requested {num_slides}")
            
            # If we got fewer slides than requested, add content slides before conclusion
            while len(slides_data) < num_slides:
                insert_position = max(2, len(slides_data) - 2)  # Insert before conclusion
                
                # Generate additional content based on slide position
                additional_topics = [
                    "Detailed Analysis and Insights",
                    "Technical Specifications and Requirements",
                    "Cost-Benefit Analysis",
                    "Stakeholder Impact Assessment",
                    "Regulatory and Compliance Considerations",
                    "Technology Integration and Innovation",
                    "Performance Metrics and Evaluation",
                    "Sustainability and Environmental Impact",
                    "Global Perspectives and Comparisons",
                    "Future Research and Development",
                    "Partnership and Collaboration Opportunities",
                    "Training and Development Needs"
                ]
                
                topic_index = (len(slides_data) - 3) % len(additional_topics)
                topic_title = additional_topics[topic_index]
                
                new_slide = {
                    "title": f"{topic_title}",
                    "content": [
                        f"Comprehensive analysis of {topic_title.lower()} in relation to {prompt_topic}",
                        f"Key factors and considerations for {topic_title.lower()}",
                        f"Strategic implications and recommendations",
                        f"Best practices and implementation guidelines"
                    ]
                }
                
                slides_data.insert(insert_position, new_slide)
            
            # If we got more slides than requested, trim excess content slides
            if len(slides_data) > num_slides:
                # Keep title, conclusion, and references, trim middle content
                title_slide = slides_data[0]
                content_slides = slides_data[1:-2][:num_slides-3]
                conclusion_slide = slides_data[-2] if len(slides_data) >= 2 else {
                    "title": "Conclusion and Next Steps",
                    "content": [
                        f"Summary: {prompt_topic} presents significant opportunities and considerations",
                        "Key takeaways from our comprehensive analysis",
                        "Strategic recommendations for implementation",
                        "Immediate action items and priorities",
                        "Long-term vision and goals",
                        "Success metrics and monitoring approach"
                    ]
                }
                references_slide = slides_data[-1] if len(slides_data) >= 1 else {
                    "title": "References and Sources",
                    "content": [
                        "1. Industry reports and publications",
                        "2. Academic research and studies",
                        "3. Expert analysis and insights"
                    ]
                }
                
                slides_data = [title_slide] + content_slides + [conclusion_slide, references_slide]

        # Ensure the last two slides are always Conclusion and References
        if num_slides >= 2:
            # Update second-to-last slide to be conclusion
            if len(slides_data) >= 2:
                conclusion_keywords = ["conclusion", "summary", "next steps", "wrap up", "final"]
                if not any(keyword in slides_data[-2]["title"].lower() for keyword in conclusion_keywords):
                    slides_data[-2] = {
                        "title": "Conclusion and Next Steps",
                        "content": [
                            f"Summary: {prompt_topic} presents significant opportunities and considerations",
                            "Key takeaways from our comprehensive analysis",
                            "Strategic recommendations for implementation",
                            "Immediate action items and priorities",
                            "Long-term vision and goals",
                            "Success metrics and monitoring approach"
                        ]
                    }
            
            # Update last slide to be references
            if len(slides_data) >= 1:
                reference_keywords = ["reference", "source", "bibliography", "citation"]
                if not any(keyword in slides_data[-1]["title"].lower() for keyword in reference_keywords):
                    slides_data[-1] = {
                        "title": "References and Sources",
                        "content": [
                            "1. Industry Research Reports and Market Analysis",
                            "2. Academic Journals and Peer-Reviewed Studies",
                            "3. Government Publications and Regulatory Documents",
                            "4. Professional Association Guidelines and Standards",
                            "5. Expert Interviews and Industry Surveys",
                            "6. Company Reports and Case Study Documentation"
                        ]
                    }

        # Final validation - ensure all slides have substantial content
        for i, slide in enumerate(slides_data):
            if not slide.get("content") or len(slide["content"]) == 0:
                slide["content"] = [
                    f"Detailed information about {slide.get('title', 'this topic')}",
                    "Key insights and analysis",
                    "Important considerations and implications",
                    "Strategic recommendations"
                ]
            elif len(slide["content"]) < 2:
                # Ensure each slide has at least 2-3 content points
                slide["content"].extend([
                    "Additional insights and analysis",
                    "Strategic implications and recommendations"
                ])

        # After successful slide generation, store presentation metadata and slides in Firestore
        if user_id:
            doc = firestore_db.collection('presentations').document()
            doc.set({
                'user_id': user_id,
                'title': prompt_topic,
                'template': template,
                'slides': slides_data,
                'created_at': firestore.SERVER_TIMESTAMP,
                'updated_at': firestore.SERVER_TIMESTAMP
            })
            # Update analytics after successful slide generation and saving
            update_analytics_on_slide(user_id, topic=prompt_topic)

        return jsonify({"slides": slides_data})

    except json.JSONDecodeError as e:
        print(f"JSON Decode Error: {e}")
        return jsonify({"error": "JSON decode error."}), 500
    except Exception as e:
        print(f"Error during slide generation: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": "Internal server error."}), 500


def apply_markdown_formatting(run, text):
    import re
    # Bold
    bold_match = re.search(r"\*\*(.*?)\*\*", text)
    if bold_match:
        run.bold = True
        text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    # Italic
    italic_match = re.search(r"\*(.*?)\*", text)
    if italic_match:
        run.italic = True
        text = re.sub(r"\*(.*?)\*", r"\1", text)
    # Underline
    underline_match = re.search(r"__(.*?)__", text)
    if underline_match:
        run.underline = True
        text = re.sub(r"__(.*?)__", r"\1", text)
    run.text = text

def download_image(url, filename):
    response = requests.get(url)
    if response.status_code == 200:
        with open(filename, "wb") as f:
            f.write(response.content)
        return filename
    return None

def apply_template_to_slide(slide, template, slide_data, ppt=None):
    # Set background color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*template["background_color"])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = slide_data.get("title", "Untitled Slide")
        p = title_shape.text_frame.paragraphs[0]
        font = p.font
        font.name = template["title_font"]["name"]
        font.size = Pt(template["title_font"]["size"])
        font.color.rgb = RGBColor(*template["title_font"]["color"])
        font.bold = template.get("title_bold", False)

    content_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if content_shape:
        content_shape.text_frame.clear()
        content_list = slide_data.get("content", [])

        # --- AUTO FONT SIZE LOGIC ---
        base_font_size = template["content_font"]["size"]
        font_size = base_font_size
        if template.get("auto_font_size", False):
            # Reduce font size if content is long (e.g. > 6 lines)
            n_items = len(content_list) if isinstance(content_list, list) else 1
            if n_items > 12:
                font_size = max(14, base_font_size - 8)
            elif n_items > 9:
                font_size = max(16, base_font_size - 6)
            elif n_items > 6:
                font_size = max(18, base_font_size - 4)
            else:
                font_size = base_font_size

        for idx, item in enumerate(content_list):
            p = content_shape.text_frame.add_paragraph()
            p.text = item
            font = p.font
            font.name = template["content_font"]["name"]
            font.size = Pt(font_size)
            font.color.rgb = RGBColor(*template["content_font"]["color"])
            font.bold = template.get("content_bold", False)
            p.level = 0 if idx == 0 else 1
    else:
        # fallback: add a textbox if placeholder not found
        left = Pt(50)
        top = Pt(150)
        width = Pt(600)
        height = Pt(350)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.clear()
        content_list = slide_data.get("content", [])
        base_font_size = template["content_font"]["size"]
        font_size = base_font_size
        if template.get("auto_font_size", False):
            n_items = len(content_list) if isinstance(content_list, list) else 1
            if n_items > 12:
                font_size = max(14, base_font_size - 8)
            elif n_items > 9:
                font_size = max(16, base_font_size - 6)
            elif n_items > 6:
                font_size = max(18, base_font_size - 4)
            else:
                font_size = base_font_size
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            font = p.font
            font.name = template["content_font"]["name"]
            font.size = Pt(font_size)
            font.color.rgb = RGBColor(*template["content_font"]["color"])
            font.bold = template.get("content_bold", False)

    if slide_data.get("image_url"):
        try:
            image_url = slide_data["image_url"]
            img_data = requests.get(image_url).content
            img_path = f"/tmp/{uuid.uuid4().hex}.png"
            with open(img_path, "wb") as f:
                f.write(img_data)
            # Place image on right half of slide, filling it edge-to-edge
            if ppt is not None:
                slide_width = ppt.slide_width
                slide_height = ppt.slide_height
                # Right half
                left = slide_width // 2
                top = 0
                width = slide_width // 2
                height = slide_height
                # Open image and crop/resize to fit right half
                with Image.open(img_path) as im:
                    # Calculate aspect ratios
                    target_ratio = width / height
                    img_ratio = im.width / im.height
                    # Crop image to match target aspect ratio (center crop)
                    if img_ratio > target_ratio:
                        # Image is wider than target: crop sides
                        new_width = int(im.height * target_ratio)
                        offset = (im.width - new_width) // 2
                        box = (offset, 0, offset + new_width, im.height)
                        im_cropped = im.crop(box)
                    else:
                        # Image is taller than target: crop top/bottom
                        new_height = int(im.width / target_ratio)
                        offset = (im.height - new_height) // 2
                        box = (0, offset, im.width, offset + new_height)
                        im_cropped = im.crop(box)
                    # Save cropped image
                    im_cropped.save(img_path)
                slide.shapes.add_picture(img_path, left, top, width, height)
            else:
                # fallback values if ppt is not provided
                left = Pt(400)
                top = Pt(0)
                width = Pt(300)
                height = Pt(225)
                slide.shapes.add_picture(img_path, left, top, width, height)
            os.remove(img_path)
        except Exception as e:
            print(f"Failed to add image to slide: {e}")


@main.route("/generate-presentation", methods=["POST", "OPTIONS"])
def generate_presentation():
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'ok'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
        response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
        return response
        
    # Define slide and editor dimensions (in inches and pixels)
    PPTX_SLIDE_WIDTH_INCHES = 13.33
    PPTX_SLIDE_HEIGHT_INCHES = 7.5
    EDITOR_SLIDE_WIDTH_PX = 1280 # Assuming this is the canvas width in the frontend editor
    EDITOR_SLIDE_HEIGHT_PX = 720 # Assuming this is the canvas height in the frontend editor

    data = request.json
    if not data:
        return jsonify({"error": "No data provided"}), 400
        
    slides_data = data.get("slides")
    if not slides_data or not isinstance(slides_data, list):
        return jsonify({"error": "Invalid slides data"}), 400
    
    global_template_id = data.get("templateId") or data.get("template")
    
    # Ensure global template ID is a string
    if isinstance(global_template_id, dict):
        global_template_id = global_template_id.get("id") or global_template_id.get("templateId")
    
    current_app.logger.info(f"Using global template ID: {global_template_id}")

    prs = PptxPresentation()
    prs.slide_width = Inches(PPTX_SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(PPTX_SLIDE_HEIGHT_INCHES)

    # Calculate conversion factors once
    px_to_in_x = PPTX_SLIDE_WIDTH_INCHES / EDITOR_SLIDE_WIDTH_PX
    px_to_in_y = PPTX_SLIDE_HEIGHT_INCHES / EDITOR_SLIDE_HEIGHT_PX

    # Helper function (can be outside or nested if only used here)
    def parse_span_style(style_str):
        styles = {}
        if not style_str:
            return styles
        for part in style_str.split(';'):
            if ':' in part:
                key, value = part.split(':', 1)
                key = key.strip().replace('-', '').lower()
                value = value.strip()
                styles[key] = value
        return styles    
    for slide_index, slide_item_data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[6]  # BLANK layout
        ppt_slide = prs.slides.add_slide(slide_layout)        # Determine template for this slide
        slide_template_id = slide_item_data.get("templateId") or slide_item_data.get("template") or global_template_id
        
        # Ensure template_id is a string, not a dict
        if isinstance(slide_template_id, dict):
            slide_template_id = slide_template_id.get("id") or slide_template_id.get("templateId")
        
        current_app.logger.debug(f"Slide {slide_index}: using template '{slide_template_id}'")
          # Apply template background if available
        template_applied = False
        if slide_template_id and isinstance(slide_template_id, str):
            # Create a simple template definition object
            template_def = {"id": slide_template_id}
              # Determine slide type based on slide position
            # Slide 1 = title, Slides 2+ = content for ALL templates
            slide_type = "title" if slide_index == 0 else "content"
            current_app.logger.debug(f"Template '{slide_template_id}' - Slide {slide_index} set to type: {slide_type}")
                
            template_applied = draw_template_slide_background(
                ppt_slide, template_def, slide_type, 
                PPTX_SLIDE_WIDTH_INCHES, PPTX_SLIDE_HEIGHT_INCHES
            )
            
            if template_applied:
                current_app.logger.info(f"Successfully applied template '{slide_template_id}' (type: {slide_type}) to slide {slide_index}")
            else:
                current_app.logger.warning(f"Failed to apply template '{slide_template_id}' to slide {slide_index}")
        else:
            if slide_template_id:
                current_app.logger.warning(f"Template '{slide_template_id}' not found or invalid for slide {slide_index}")
            else:
                current_app.logger.debug(f"No template specified for slide {slide_index}")        # Fallback: apply simple background if no template was applied
        if not template_applied:
            background_data = slide_item_data.get("background", {})
            bg_fill_hex = background_data.get("fill", "#FFFFFF")
            if bg_fill_hex:
                try:
                    fill = ppt_slide.background.fill
                    fill.solid()
                    rgb_tuple = hex_to_rgb(bg_fill_hex)
                    fill.fore_color.rgb = RGBColor(*rgb_tuple)
                except Exception as e:
                    current_app.logger.error(f"Error setting background color: {e}", exc_info=True)


        # Collect elements to be rendered based on zIndex
        elements_to_render = []

        textboxes_data = slide_item_data.get("textboxes", [])
        for tb_data in textboxes_data:
            elements_to_render.append({
                "type": "textbox",
                "data": tb_data,
                "zIndex": int(tb_data.get("zIndex", 100)) # Default zIndex for textboxes
            })

        images_data = slide_item_data.get("images", []) # Process multiple images
        for img_data in images_data:
            elements_to_render.append({
                "type": "image",
                "data": img_data,
                "zIndex": int(img_data.get("zIndex", 101)) # Images have zIndex
            })
          # Sort elements by zIndex: lower zIndex elements are added first (appear "behind")
        elements_to_render.sort(key=lambda el: el["zIndex"])
          # Render elements in sorted order
        for element in elements_to_render:
            el_data = element["data"]
            el_type = element["type"]
            
            if el_type == "textbox":
                try:
                    x_px = float(el_data.get("x", 0))
                    y_px = float(el_data.get("y", 0))
                    width_px = float(el_data.get("width", 100))
                    height_px = float(el_data.get("height", 50))
                    left = Inches(x_px * px_to_in_x)
                    top = Inches(y_px * px_to_in_y)
                    width = Inches(width_px * px_to_in_x)
                    height = Inches(height_px * px_to_in_y)                  
                    if width <= Inches(0) or height <= Inches(0):
                        current_app.logger.warning(f"Skipping textbox with invalid dimensions: w_px={width_px}, h_px={height_px}")
                        continue

                    shape = ppt_slide.shapes.add_textbox(left, top, width, height)
                    tf = shape.text_frame
                    tf.word_wrap = True 
                    tf.auto_size = MSO_AUTO_SIZE.NONE # Use explicit height
                    tf.margin_bottom = Inches(0.05) 
                    tf.margin_left = Inches(0.1)
                    tf.margin_right = Inches(0.1)
                    tf.margin_top = Inches(0.05)
                    tf.clear()

                    text_content = el_data.get("text", "")
                    default_font_family = el_data.get("fontFamily", "Arial")
                    default_font_size_pt = float(el_data.get("fontSize", 18))
                    default_font_color_hex = el_data.get("fill", "#000000")
                    default_font_color_rgb = hex_to_rgb(default_font_color_hex)
                    
                    default_font_style_data = el_data.get("fontStyle", {})
                    default_bold = default_font_style_data.get("bold", False)
                    default_italic = default_font_style_data.get("italic", False)
                    default_underline = default_font_style_data.get("underline", False)
                    
                    default_align_str = el_data.get("align", "left").upper()
                    align_map = {
                        "LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER,
                        "RIGHT": PP_ALIGN.RIGHT, "JUSTIFY": PP_ALIGN.JUSTIFY,
                    }
                    default_alignment = align_map.get(default_align_str, PP_ALIGN.LEFT)
                    
                    default_line_height_multiplier = el_data.get("lineHeight") # e.g., 1, 1.15, 1.5
                    default_paragraph_spacing_pt = float(el_data.get("paragraphSpacing", 0))
                    is_bulleted = el_data.get("bullets", False)

                    paragraphs_text = text_content.split('\\\\n') # Split by literal \\n from JSON
                    
                    if not paragraphs_text and not text_content.strip(): # Handle completely empty textbox or textbox with only &nbsp;
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = " " # Add a space to make it selectable and visible if it has dimensions
                        run.font.size = Pt(1)
                        continue # Move to next element

                    for para_idx, para_text_html in enumerate(paragraphs_text):
                        p = tf.add_paragraph()
                        p.alignment = default_alignment
                        if default_line_height_multiplier and isinstance(default_line_height_multiplier, (int, float)):
                            try:
                                p.line_spacing = float(default_line_height_multiplier)
                            except ValueError:
                                current_app.logger.warning(f"Invalid line height value: {default_line_height_multiplier}")
                        
                        if para_idx > 0 and default_paragraph_spacing_pt > 0:
                             p.space_before = Pt(default_paragraph_spacing_pt)
                        
                        if is_bulleted and para_text_html.strip(): # Add bullet only if line has content
                            p.level = 0

                        # Parse HTML-like content (spans) for rich text
                        # Replace &nbsp; with space for BeautifulSoup processing
                        soup = BeautifulSoup(f"<div>{para_text_html.replace('&nbsp;', ' ')}</div>", "html.parser")
                        
                        if not soup.div.contents: # Handle paragraph that becomes empty after parsing (e.g. only &nbsp;)
                            if is_bulleted: # If it was supposed to be a bullet, keep the paragraph for the bullet point
                                pass # The paragraph is already added, bullet will show if level is set
                            elif len(paragraphs_text) > 1 or para_text_html: # Preserve empty line if it's not a truly empty single-line textbox
                                run = p.add_run()
                                run.text = " " # Add a space to make the line take height
                            continue # Next content_node or next paragraph

                        for content_node in soup.div.contents:
                            run = p.add_run()
                            text_to_add = ""
                            
                            run_font_family = default_font_family
                            run_font_size_pt = default_font_size_pt
                            run_font_color_rgb = default_font_color_rgb
                            run_bold = default_bold
                            run_italic = default_italic
                            run_underline = default_underline

                            if content_node.name == 'span':
                                text_to_add = content_node.get_text()
                                span_style_str = content_node.get('style', '')
                                span_styles = parse_span_style(span_style_str)
                                if 'fontfamily' in span_styles: run_font_family = span_styles['fontfamily']
                                if 'fontsize' in span_styles:
                                    try: run_font_size_pt = float(str(span_styles['fontsize']).replace('pt','').replace('px',''))
                                    except ValueError: pass
                                if 'color' in span_styles: 
                                    try: run_font_color_rgb = hex_to_rgb(span_styles['color'])
                                    except: pass # Ignore invalid color
                                if 'bold' in span_styles: run_bold = str(span_styles['bold']).lower() == 'true'
                                if 'italic' in span_styles: run_italic = str(span_styles['italic']).lower() == 'true'
                                if 'underline' in span_styles: run_underline = str(span_styles['underline']).lower() == 'true'
                            
                            elif content_node.name is None: # Plain text node
                                text_to_add = str(content_node)
                            
                            if text_to_add:
                                run.text = text_to_add
                                run.font.name = run_font_family
                                run.font.size = Pt(run_font_size_pt)
                                if run_font_color_rgb: 
                                    try:
                                        run.font.color.rgb = RGBColor(*run_font_color_rgb)
                                    except Exception as color_error:
                                        current_app.logger.warning(f"Failed to set font color: {color_error}")
                                run.font.bold = run_bold
                                run.font.italic = run_italic
                                run.font.underline = run_underline
                        
                        if not p.runs and not (is_bulleted and para_text_html.strip()): 
                             if len(paragraphs_text) > 1 or para_text_html: 
                                 run = p.add_run()
                                 run.text = " " 

                    if not tf.paragraphs: # Final check if textbox ended up with no paragraphs at all
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = " " 
                        run.font.size = Pt(1)
                except Exception as e:
                    tb_id_log = "Unknown Textbox"
                    if isinstance(el_data, dict):
                        tb_id_log = el_data.get('id', 'N/A')
                    current_app.logger.error(f"Error processing textbox: {tb_id_log}. Error: {e}", exc_info=True)

            elif el_type == "image":
                try:
                    img_src_base64 = el_data.get("src")
                    if img_src_base64 and img_src_base64.startswith('data:image'):
                        header, encoded = img_src_base64.split(',', 1)
                        img_bytes = base64.b64decode(encoded)
                        img_stream = BytesIO(img_bytes)

                        img_x_px = float(el_data.get("x", 0))
                        img_y_px = float(el_data.get("y", 0))
                        img_width_px = float(el_data.get("width", 100))
                        img_height_px = float(el_data.get("height", 100))
                        if img_width_px <= 0 or img_height_px <= 0:
                            current_app.logger.warning(f"Skipping image with zero/negative pixel dimensions: w={img_width_px}, h={img_height_px}")
                            continue                       
                        img_left = Inches(img_x_px * px_to_in_x)
                        img_top = Inches(img_y_px * px_to_in_y)
                        img_width = Inches(img_width_px * px_to_in_x)
                        img_height = Inches(img_height_px * px_to_in_y)
                        
                        if img_width <= Inches(0) or img_height <= Inches(0):
                            current_app.logger.warning(f"Skipping image with zero/negative inch dimensions after conversion: w_in={img_width}, h_in={img_height}")
                            continue

                        ppt_slide.shapes.add_picture(img_stream, img_left, img_top, width=img_width, height=img_height)
                except Exception as e:
                    img_id_log = "Unknown Image"
                    if isinstance(el_data, dict):
                        img_id_log = el_data.get('id', 'N/A')
                    current_app.logger.error(f"Error processing image: {img_id_log}. Error: {e}", exc_info=True)

    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    return send_file(
        file_stream,
        as_attachment=True,
        download_name="smartslide_presentation.pptx", # Changed name slightly
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def draw_template_slide_background(slide, template_def, slide_type="default", slide_width_inches=13.33, slide_height_inches=7.5):
    try:
        from pptx.util import Inches
        slide_width = Inches(slide_width_inches)
        slide_height = Inches(slide_height_inches)
        template_id = template_def.get("id") or template_def.get("name", "").lower()

        # --- BUSINESS TEMPLATE ---
        if template_id == "tailwind-business":
            img_path = os.path.join(
                os.path.dirname(__file__),
                "static", "template_backgrounds",
                "tailwind-business_title.png" if slide_type == "title" else "tailwind-business_content.png"
            )
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)
                return True        # --- CREATIVE TEMPLATE ---
        elif template_id == "tailwind-creative":
            img_path = os.path.join(
                os.path.dirname(__file__),
                "static", "template_backgrounds",
                "tailwind-creative_title.png" if slide_type == "title" else "tailwind-creative_content.png"
            )
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)
                return True

        # --- EDUCATION TEMPLATE ---
        elif template_id == "tailwind-education":
            img_path = os.path.join(
                os.path.dirname(__file__),
                "static", "template_backgrounds",
                "tailwind-education_title.png" if slide_type == "title" else "tailwind-education_content.png"
            )
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)
                return True        # --- ABSTRACT GRADIENT TEMPLATE ---
        elif template_id == "tailwind-abstract-gradient":
            img_path = os.path.join(
                os.path.dirname(__file__),
                "static", "template_backgrounds",
                "tailwind-abstract-gradient_title.png" if slide_type == "title" else "tailwind-abstract-gradient_content.png"
            )
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)
                return True

        # No fallback: only use image backgrounds for all templates
        return False
    except Exception as e:
        current_app.logger.error(f"Error drawing template background: {e}")
        return False

# --- FIREBASE ANALYTICS ENDPOINT ---
@main.route('/analytics/<user_id>', methods=['GET', 'OPTIONS'])
def get_user_analytics(user_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    # Get analytics document for the user
    analytics_doc = firestore_db.collection('analytics').document(str(user_id)).get()
    analytics_data = analytics_doc.to_dict() if analytics_doc.exists else {}

    # Get presentations for user
    presentations_ref = firestore_db.collection('presentations')
    presentations = [doc.to_dict() for doc in presentations_ref.where('user_id', '==', user_id).stream()]

    # Slides created over time (monthly)
    from collections import Counter
    monthly_slides_counter = Counter()
    for p in presentations:
        created_at = p.get('created_at')
        if created_at:
            # Firestore timestamp to ISO string, then to YYYY-MM
            if hasattr(created_at, 'isoformat'):
                month_year = created_at.isoformat()[:7]
            else:
                month_year = str(created_at)[:7]
            monthly_slides_counter[month_year] += 1

    # Most common topics (using presentation titles)
    topic_counts_counter = Counter()
    for p in presentations:
        title = p.get('title')
        if title:
            topic_counts_counter[title.lower()] += 1

    slides_generated = analytics_data.get('slides_created', 0)
    quizzes_generated = len(list(firestore_db.collection('saved_quizzes').where('user_id', '==', user_id).stream()))
    scripts_generated = len(list(firestore_db.collection('saved_scripts').where('user_id', '==', user_id).stream()))
    last_active = analytics_data.get('last_generated_at')
    last_active_iso = last_active.isoformat() if last_active and hasattr(last_active, 'isoformat') else None

    return jsonify({
        "monthly": dict(monthly_slides_counter),
        "topics": dict(topic_counts_counter),
        "slides_generated": slides_generated,
        "quizzes_generated": quizzes_generated,
        "scripts_generated": scripts_generated,
        "last_active": last_active_iso
    }), 200



# Update the paste-and-create endpoint around line 1522
@main.route('/paste-and-create', methods=['POST', 'OPTIONS'])
def paste_and_create():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    data = request.get_json()
    pasted_text = data.get('text', '').strip()
    language = data.get('language', 'English')
    
    # Fix: Support up to 30 slides and validate properly
    try:
        num_slides = int(data.get('numSlides', 6))
        if num_slides <= 0 or num_slides > 30:
            return jsonify({'error': 'Number of slides must be between 1 and 30'}), 400
    except (ValueError, TypeError):
        return jsonify({'error': 'Invalid number of slides'}), 400

    if not pasted_text:
        return jsonify({'error': 'No text provided'}), 400

    # Extract topic/title
    lines = [line.strip() for line in pasted_text.split('\n') if line.strip()]
    topic = lines[0] if lines else "Untitled Topic"

    # Structure the text for slides
    prompt = f"""
You are an assistant that structures pasted text into a professional presentation outline.
Given the following text, extract the main topic as the title, and organize the content into exactly {num_slides} slides.

Requirements:
- Each slide should have a concise and clear title and substantial content.
- Format the content as bullet points or short paragraphs.
- Use Markdown for formatting: **bold** for emphasis, *italic* for highlights, and __underline__ for key terms.
- If the text is long, summarize and split it logically across slides.
- Ensure ALL {num_slides} slides have meaningful content.
- The last two slides must be:
  - Slide {num_slides-1}: Conclusion and Next Steps
  - Slide {num_slides}: References and Sources

Structure:
- Slide 1: Title slide with the topic and description
- Slides 2-{num_slides-2}: Main content slides covering different aspects
- Slide {num_slides-1}: Conclusion and Next Steps
- Slide {num_slides}: References and Sources

Text:
---
{pasted_text}
---

Output a JSON array where each object has:
- 'title': string (slide title)
- 'content': array of strings (slide content, Markdown allowed)
- 'image_prompt': string (optional, for relevant image description)

Make sure to generate exactly {num_slides} slides with substantial content for each.
"""

    try:
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        
        # Increase token limit for larger presentations
        max_tokens = 4096
        if num_slides > 15:
            max_tokens = 6144
        elif num_slides > 25:
            max_tokens = 8192
            
        payload = {
            "model": "llama3-8b-8192",
            "messages": [
                {"role": "system", "content": "You are a helpful assistant that generates slide content in JSON format. Always ensure comprehensive content and that the last two slides are Conclusion and References."},
                {"role": "user", "content": prompt}
            ],
            "max_tokens": max_tokens,
            "temperature": 0.4
        }
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        result = response.json()
        model_output = result["choices"][0]["message"]["content"]

        # Extract JSON array from the response
        import re, json
        match = re.search(r'\[\s*{.*}\s*\]', model_output, re.DOTALL)
        if not match:
            return jsonify({"error": "Failed to parse slides output."}), 500
        slides_data = json.loads(match.group(0))

        # Validate and adjust slide count with proper conclusion and references
        if len(slides_data) != num_slides:
            while len(slides_data) < num_slides:
                slides_data.append({
                    "title": f"Additional Content {len(slides_data) + 1}",
                    "content": ["Additional content will be added here."]
                })
            slides_data = slides_data[:num_slides]

        # Ensure last two slides are conclusion and references
        if num_slides >= 2:
            slides_data[-2] = {
                "title": "Conclusion and Next Steps",
                "content": [
                    "Summary of key findings and insights",
                    "Main takeaways from the analysis",
                    "Recommended next steps and actions",
                    "Future considerations and planning"
                ]
            }
            slides_data[-1] = {
                "title": "References and Sources",
                "content": [
                    "1. Original document analysis",
                    "2. Supporting research and data",
                    "3. Industry best practices",
                    "4. Expert recommendations"
                ]
            }

        # Ensure all slides have substantial content
        for slide in slides_data:
            if not slide.get("content") or len(slide["content"]) < 2:
                slide["content"] = [
                    f"Key information about {slide.get('title', 'this topic')}",
                    "Important insights and analysis",
                    "Strategic considerations and recommendations"
                ]

        detected_topic = slides_data[0]["title"] if slides_data and "title" in slides_data[0] else topic

        return jsonify({
            "slides": slides_data,
            "topic": detected_topic
        }), 200

    except Exception as e:
        current_app.logger.error(f"Error in /upload-file: {e}", exc_info=True)
        return jsonify({"error": "Failed to generate slides from file."}), 500
    
# --- EXPORT QUIZ AND SCRIPT TO WORD (STUBS) ---
@main.route('/export-quiz-word', methods=['POST', 'OPTIONS'])
def export_quiz_word():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'No data provided'}), 400
            
        quiz = data.get('quiz')
        if not quiz:
            return jsonify({'error': 'No quiz data provided'}), 400
            
        doc = Document()
        doc.add_heading('Generated Quiz', 0)
        
        for idx, q in enumerate(quiz, 1):
            question = q.get('question', '')
            choices = q.get('choices', [])
            answer = q.get('answer', '')
            
            doc.add_paragraph(f"{idx}. {question}", style='List Number')
            
            if choices:
                for cidx, choice in enumerate(choices, 1):
                    doc.add_paragraph(f"    {chr(64+cidx)}. {choice}")
            
            doc.add_paragraph(f"    Answer: {answer}")
            doc.add_paragraph("")  # Add spacing
        
        f = BytesIO()
        doc.save(f)
        f.seek(0)
        
        return send_file(
            f, 
            as_attachment=True, 
            download_name="generated_quiz.docx", 
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
    except Exception as e:
        current_app.logger.error(f"Error in export-quiz-word: {e}", exc_info=True)
        return jsonify({'error': 'Failed to export quiz'}), 500
@main.route('/export-script-word', methods=['POST', 'OPTIONS'])
def export_script_word():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    data = request.get_json()
    script = data.get('script')
    if not script:
        return jsonify({'error': 'No script data provided'}), 400
    doc = Document()
    doc.add_heading('Speaker Script', 0)
    doc.add_paragraph(script)
    f = BytesIO()
    doc.save(f)
    f.seek(0)
    return send_file(f, as_attachment=True, download_name="script.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

@main.route('/export-quiz/<quiz_id>/word', methods=['GET', 'OPTIONS'])
def export_quiz_by_id_word(quiz_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    quiz_doc = firestore_db.collection('saved_quizzes').document(str(quiz_id)).get()
    if not quiz_doc.exists:
        return jsonify({'error': 'Quiz not found'}), 404
    doc_data = quiz_doc.to_dict()
    # Try both 'quiz' and 'content' fields
    quiz = doc_data.get('quiz') or doc_data.get('content')
    if isinstance(quiz, str):
        try:
            quiz = json.loads(quiz)
        except Exception:
            pass
    if not quiz or not isinstance(quiz, list):
        return jsonify({'error': 'Quiz data missing or invalid'}), 400
    doc = Document()
    doc.add_heading('Quiz', 0)
    for idx, q in enumerate(quiz, 1):
        question = q.get('question', '')
        choices = q.get('choices', [])
        answer = q.get('answer', '')
        doc.add_paragraph(f"{idx}. {question}", style='List Number')
        if choices:
            for cidx, choice in enumerate(choices, 1):
                doc.add_paragraph(f"    {chr(64+cidx)}. {choice}")
        doc.add_paragraph(f"    Answer: {answer}")
    f = BytesIO()
    doc.save(f)
    f.seek(0)
    return send_file(f, as_attachment=True, download_name=f"quiz_{quiz_id}.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

@main.route('/export-script/<script_id>/word', methods=['GET', 'OPTIONS'])
def export_script_by_id_word(script_id):
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    script_doc = firestore_db.collection('saved_scripts').document(str(script_id)).get()
    if not script_doc.exists:
        return jsonify({'error': 'Script not found'}), 404
    doc_data = script_doc.to_dict()
    # Try both 'script' and 'content' fields
    script = doc_data.get('script') or doc_data.get('content')
    if isinstance(script, dict) and 'content' in script:
        script = script['content']
    if not script or not isinstance(script, str):
        return jsonify({'error': 'Script data missing or invalid'}), 400
    doc = Document()
    doc.add_heading('Speaker Script', 0)
    doc.add_paragraph(script)
    f = BytesIO()
    doc.save(f)
    f.seek(0)
    return send_file(f, as_attachment=True, download_name=f"script_{script_id}.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# --- FORGOT PASSWORD ENDPOINTS ---
@main.route('/forgot-password/send-verification', methods=['POST', 'OPTIONS'])
def send_password_reset_verification():
    """Send password reset verification code via email or SMS"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    data = request.get_json()
    identifier = data.get('identifier')  # email or phone number
    method = data.get('method')  # 'email' or 'phone'
    check_only = data.get('check_only', False)  # Just check if user exists
    
    if not identifier or not method:
        return jsonify({'error': 'Identifier and method are required'}), 400
    
    if method not in ['email', 'phone']:
        return jsonify({'error': 'Method must be email or phone'}), 400
    
    try:
        users_ref = firestore_db.collection('users')
        
        if method == 'email':
            # Find user by email
            users_query = users_ref.where('email', '==', identifier).get()
            if not users_query:
                return jsonify({'error': 'No user found with this email'}), 404
            
            # If just checking, return success
            if check_only:
                return jsonify({'message': 'User found'}), 200
            
            # Generate 6-digit verification code
            import secrets
            import string
            verification_code = ''.join(secrets.choice(string.digits) for _ in range(6))
            
            user_doc = users_query[0]
            user_id = user_doc.id
            
            # Store verification code in Firestore with expiration
            from datetime import datetime, timedelta
            reset_ref = firestore_db.collection('password_resets').document()
            reset_ref.set({
                'user_id': user_id,
                'email': identifier,
                'verification_code': verification_code,
                'expires_at': datetime.utcnow().replace(tzinfo=timezone.utc) + timedelta(minutes=10),
                'created_at': firestore.SERVER_TIMESTAMP,
                'verified': False
            })
            
            # Send email verification
            success, message = send_email_verification(identifier, verification_code)
            
            if success:
                return jsonify({
                    'message': 'Verification code sent to your email',
                    'reset_id': reset_ref.id
                }), 200
            else:
                # Fallback: log the code for development
                current_app.logger.info(f"Email service failed, password reset code for {identifier}: {verification_code}")
                return jsonify({
                    'message': 'Verification code sent to your email',
                    'reset_id': reset_ref.id,
                    # Remove this line in production:
                    'verification_code': verification_code
                }), 200
        elif method == 'phone':
            # Find user by contact number
            users_query = users_ref.where('contact_number', '==', identifier).get()
            if not users_query:
                return jsonify({'error': 'No user found with this phone number'}), 404
            
            # If just checking, return success
            if check_only:
                return jsonify({'message': 'User found'}), 200
            
            # Generate verification code for SMS
            import secrets
            import string
            verification_code = ''.join(secrets.choice(string.digits) for _ in range(6))
            
            user_doc = users_query[0]
            user_id = user_doc.id
            
            # Store verification code in Firestore
            from datetime import datetime, timedelta
            reset_ref = firestore_db.collection('password_resets').document()
            reset_ref.set({
                'user_id': user_id,
                'phone': identifier,
                'verification_code': verification_code,
                'expires_at': datetime.utcnow().replace(tzinfo=timezone.utc) + timedelta(minutes=10),
                'created_at': firestore.SERVER_TIMESTAMP,
                'verified': False
            })
            
            # Send SMS verification
            success, message = send_sms_verification(identifier, verification_code)
            
            if success:
                return jsonify({
                    'message': 'Verification code sent to your phone',
                    'reset_id': reset_ref.id
                }), 200
            else:
                # Fallback: log the code for development
                current_app.logger.info(f"SMS service failed, password reset code for {identifier}: {verification_code}")
                return jsonify({
                    'message': 'Verification code sent to your phone',
                    'reset_id': reset_ref.id,
                    # Remove this line in production:
                    'verification_code': verification_code
                }), 200
            
    except Exception as e:
        current_app.logger.error(f"Error in password reset verification: {e}")
        return jsonify({'error': 'Failed to send verification code'}), 500

@main.route('/forgot-password/verify-code', methods=['POST', 'OPTIONS'])
def verify_password_reset_code():
    """Verify the password reset code"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.get_json()
    reset_id = data.get('reset_id')
    verification_code = data.get('verification_code')
    
    if not reset_id or not verification_code:
        return jsonify({'error': 'Reset ID and verification code are required'}), 400
    
    try:
        # Get the password reset document
        reset_ref = firestore_db.collection('password_resets').document(reset_id)
        reset_doc = reset_ref.get()
        
        if not reset_doc.exists:
            return jsonify({'error': 'Invalid reset request'}), 404
        
        reset_data = reset_doc.to_dict()
        
        # Check if already verified
        if reset_data.get('verified'):
            return jsonify({'error': 'Verification code already used'}), 400
          # Check expiration - handle timezone-aware comparison
        from datetime import datetime, timezone
        expires_at = reset_data['expires_at']
        
        # Convert expires_at to timezone-aware datetime properly
        if hasattr(expires_at, 'timestamp'):
            # If it's a Firestore timestamp, convert to datetime with UTC timezone
            expires_at = datetime.fromtimestamp(expires_at.timestamp(), tz=timezone.utc)
        elif hasattr(expires_at, 'tzinfo') and expires_at.tzinfo is None:
            # If it's a naive datetime, assume UTC
            expires_at = expires_at.replace(tzinfo=timezone.utc)
        elif not hasattr(expires_at, 'tzinfo'):
            # If it's a datetime without timezone info, assume UTC
            expires_at = expires_at.replace(tzinfo=timezone.utc)
        
        current_time = datetime.now(timezone.utc)
        
        if current_time > expires_at:
            return jsonify({'error': 'Verification code expired'}), 400
        
        # Verify the code
        if reset_data['verification_code'] != verification_code:
            return jsonify({'error': 'Invalid verification code'}), 400
        
        # Mark as verified
        reset_ref.update({
            'verified': True,
            'verified_at': firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            'message': 'Verification code verified successfully',
            'reset_id': reset_id
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Error verifying reset code: {e}")
        return jsonify({'error': 'Failed to verify code'}), 500

@main.route('/forgot-password/reset-password', methods=['POST', 'OPTIONS'])
def reset_password():
    """Reset the user's password after verification"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.get_json()
    reset_id = data.get('reset_id')
    new_password = data.get('new_password')
    
    if not reset_id or not new_password:
        return jsonify({'error': 'Reset ID and new password are required'}), 400
    
    try:
        # Get the password reset document
        reset_ref = firestore_db.collection('password_resets').document(reset_id)
        reset_doc = reset_ref.get()
        
        if not reset_doc.exists:
            return jsonify({'error': 'Invalid reset request'}), 404
        
        reset_data = reset_doc.to_dict()
        
        # Check if verified
        if not reset_data.get('verified'):
            return jsonify({'error': 'Verification required before password reset'}), 400          # Check expiration (allow 30 minutes after verification)
        from datetime import datetime, timedelta, timezone
        verified_at = reset_data.get('verified_at')
        if verified_at:
            # Handle Firestore timestamp properly
            if hasattr(verified_at, 'timestamp'):
                verified_at = datetime.fromtimestamp(verified_at.timestamp(), tz=timezone.utc)
            elif hasattr(verified_at, 'tzinfo') and verified_at.tzinfo is None:
                verified_at = verified_at.replace(tzinfo=timezone.utc)
            elif not hasattr(verified_at, 'tzinfo'):
                verified_at = verified_at.replace(tzinfo=timezone.utc)
            
            current_time = datetime.now(timezone.utc)
            if current_time > verified_at + timedelta(minutes=30):
                return jsonify({'error': 'Reset session expired'}), 400
        
        # Update user password
        user_id = reset_data['user_id']
        users_ref = firestore_db.collection('users')
        user_ref = users_ref.document(user_id)
        
        user_ref.update({
            'password_hash': generate_password_hash(new_password),
            'updated_at': firestore.SERVER_TIMESTAMP
        })
        
        # Delete the reset document for security
        reset_ref.delete()
        
        return jsonify({
            'message': 'Password reset successfully'
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Error resetting password: {e}")
        return jsonify({'error': 'Failed to reset password'}), 500

@main.route('/forgot-password/verify-firebase-sms', methods=['POST', 'OPTIONS'])
def verify_firebase_sms():
    """Verify Firebase SMS authentication and create password reset session"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.get_json()
    identifier = data.get('identifier')  # original phone number entered
    id_token = data.get('id_token')      # Firebase ID token
    phone_number = data.get('phone_number')  # verified phone number from Firebase
    
    if not all([identifier, id_token, phone_number]):
        return jsonify({'error': 'Missing required fields'}), 400
    
    try:
        # Verify the Firebase ID token
        decoded_token = auth.verify_id_token(id_token)
        firebase_phone = decoded_token.get('phone_number')
        
        # Ensure the phone number matches
        if firebase_phone != phone_number:
            return jsonify({'error': 'Phone number verification failed'}), 400
        
        # Find user in our database by contact number
        users_ref = firestore_db.collection('users')
        users_query = users_ref.where('contact_number', '==', identifier).get()
        
        if not users_query:
            return jsonify({'error': 'No user found with this phone number'}), 404
        
        user_doc = users_query[0]
        user_id = user_doc.id
        
        # Create a password reset session (similar to email verification)
        from datetime import datetime, timedelta, timezone
        reset_ref = firestore_db.collection('password_resets').document()
        reset_ref.set({
            'user_id': user_id,
            'phone': identifier,
            'firebase_verified': True,
            'firebase_phone': firebase_phone,
            'expires_at': datetime.now(timezone.utc) + timedelta(minutes=30),
            'created_at': firestore.SERVER_TIMESTAMP,
            'verified': True,  # Already verified by Firebase
            'verified_at': firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            'message': 'Phone number verified successfully',
            'reset_id': reset_ref.id
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Error verifying Firebase SMS: {e}")
        return jsonify({'error': 'Failed to verify phone number'}), 500

# --- CUSTOM TEMPLATE UPLOAD ENDPOINTS ---

@main.route('/upload-template', methods=['POST', 'OPTIONS'])
def upload_template():
    """Upload a custom template image"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    try:
        # Check if file is in request
        if 'template' not in request.files:
            return jsonify({'error': 'No template file provided'}), 400
        
        file = request.files['template']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Get form data
        template_name = request.form.get('name', '').strip()
        user_id = request.form.get('user_id', '').strip()
        
        if not template_name or not user_id:
            return jsonify({'error': 'Template name and user ID are required'}), 400
        
        # Validate file type (PNG only)
        if not file.filename.lower().endswith('.png'):
            return jsonify({'error': 'Only PNG files are allowed'}), 400
        
        # Validate file size (max 5MB)
        file.seek(0, 2)  # Seek to end
        file_size = file.tell()
        file.seek(0)  # Reset to beginning
        
        if file_size > 5 * 1024 * 1024:  # 5MB
            return jsonify({'error': 'File size must be less than 5MB'}), 400
        
        # Generate unique filename
        import uuid
        file_extension = '.png'
        unique_filename = f"custom_template_{user_id}_{uuid.uuid4().hex[:8]}{file_extension}"
        
        # Create upload directory if it doesn't exist
        upload_dir = os.path.join(current_app.root_path, 'static', 'custom_templates')
        os.makedirs(upload_dir, exist_ok=True)
        
        # Save file
        file_path = os.path.join(upload_dir, unique_filename)
        file.save(file_path)
        
        # Store template metadata in Firestore
        template_ref = firestore_db.collection('custom_templates').document()
        template_data = {
            'user_id': user_id,
            'name': template_name,
            'filename': unique_filename,
            'file_path': f'/static/custom_templates/{unique_filename}',
            'created_at': firestore.SERVER_TIMESTAMP,
            'file_size': file_size
        }
        template_ref.set(template_data)
        
        return jsonify({
            'message': 'Template uploaded successfully',
            'template_id': template_ref.id,
            'template': {
                'id': template_ref.id,
                'name': template_name,
                'preview': f'/static/custom_templates/{unique_filename}'
            }
        }), 201
        
    except Exception as e:
        current_app.logger.error(f"Error uploading template: {e}")
        return jsonify({'error': 'Failed to upload template'}), 500

@main.route('/user-templates/<user_id>', methods=['GET', 'OPTIONS'])
def get_user_templates(user_id):
    """Get all custom templates for a user"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    try:
        # Fetch user's custom templates from Firestore
        templates_ref = firestore_db.collection('custom_templates')
        user_templates = templates_ref.where('user_id', '==', user_id).get()
        
        templates = []
        for template_doc in user_templates:
            template_data = template_doc.to_dict()
            templates.append({
                'id': template_doc.id,
                'name': template_data.get('name'),
                'preview': template_data.get('file_path'),
                'created_at': template_data.get('created_at'),
                'file_size': template_data.get('file_size', 0)
            })
        
        # Sort by creation date (newest first)
        templates.sort(key=lambda x: x.get('created_at') or '', reverse=True)
        
        return jsonify({'templates': templates}), 200
        
    except Exception as e:
        current_app.logger.error(f"Error fetching user templates: {e}")
        return jsonify({'error': 'Failed to fetch templates'}), 500

@main.route('/delete-template/<template_id>', methods=['DELETE', 'OPTIONS'])
def delete_template(template_id):
    """Delete a custom template"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    try:
        # Get template document
        template_ref = firestore_db.collection('custom_templates').document(template_id)
        template_doc = template_ref.get()
        
        if not template_doc.exists:
            return jsonify({'error': 'Template not found'}), 404
        
        template_data = template_doc.to_dict()
        filename = template_data.get('filename')
        
        # Delete file from filesystem
        if filename:
            file_path = os.path.join(current_app.root_path, 'static', 'custom_templates', filename)
            if os.path.exists(file_path):
                os.remove(file_path)
        
        # Delete document from Firestore
        template_ref.delete()
        
        return jsonify({'message': 'Template deleted successfully'}), 200
        
    except Exception as e:
        current_app.logger.error(f"Error deleting template: {e}")
        return jsonify({'error': 'Failed to delete template'}), 500

@main.route('/static/custom_templates/<filename>')
def serve_custom_template(filename):
    """Serve custom template files"""
    try:
        upload_dir = os.path.join(current_app.root_path, 'static', 'custom_templates')
        return send_from_directory(upload_dir, filename)
    except Exception as e:
        current_app.logger.error(f"Error serving custom template: {e}")
        return jsonify({'error': 'Failed to serve template file'}), 500

# --- CHATBOT ENDPOINT ---
@main.route('/chatbot', methods=['POST', 'OPTIONS'])
def chatbot():
    """Chatbot endpoint to answer questions about SmartSlide features"""
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
    
    data = request.get_json()
    user_message = data.get('message', '').strip()
    
    if not user_message:
        return jsonify({'error': 'Message is required'}), 400
    
    # SmartSlide knowledge base
    smartslide_context = """
    SmartSlide is a comprehensive presentation creation tool with the following features:

    CORE FEATURES:
    - AI-powered slide generation using advanced language models
    - Multiple presentation templates and themes
    - Slide editing with drag-and-drop functionality
    - Text formatting (fonts, colors, alignment, bullets)
    - Image integration and background customization
    - PowerPoint export functionality
    - Quiz generation from presentation content
    - Speaker script generation
    - Analytics and usage tracking

    SLIDE CREATION:
    - Generate slides from topics/prompts
    - Import content from files (PDF, DOCX, TXT, CSV, Excel)
    - Paste and create from existing text
    - Support for multiple languages
    - Customizable number of slides (1-30)
    - Professional templates with various themes

    EDITING CAPABILITIES:
    - Rich text editor with formatting options
    - Font selection from 42+ font families
    - Color customization for text and backgrounds
    - Image upload and positioning
    - Template switching and customization
    - Undo/redo functionality

    EXPORT OPTIONS:
    - Download as PowerPoint (.pptx)
    - Export quizzes to Word documents
    - Export scripts to Word documents
    - Save presentations to cloud storage

    USER FEATURES:
    - User registration and authentication
    - Dashboard with recent presentations
    - Saved quizzes and scripts management
    - User analytics and statistics
    - Custom template uploads
    - Account profile management

    ADDITIONAL TOOLS:
    - Quiz generator with multiple choice and identification questions
    - Speaker script generator for presentations
    - File upload support for various formats
    - Real-time collaboration features
    - Responsive design for mobile and desktop

    TECHNICAL CAPABILITIES:
    - Firebase integration for data storage
    - AI-powered content generation
    - Secure user authentication
    - Cloud-based presentation storage
    - Cross-platform compatibility
    """
    
    try:
        # Create a focused prompt for the chatbot
        chatbot_prompt = f"""
        You are SmartSlide Assistant, a helpful chatbot that answers questions about SmartSlide, an AI-powered presentation creation tool.

        Context about SmartSlide:
        {smartslide_context}

        User Question: {user_message}

        Please provide a helpful, accurate, and friendly response about SmartSlide's features, capabilities, or how to use the application. Keep your response concise but informative. If the question is not related to SmartSlide, politely redirect the conversation back to SmartSlide features.

        Response:
        """
        
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": "llama3-8b-8192",
            "messages": [
                {"role": "system", "content": "You are SmartSlide Assistant, a helpful chatbot for the SmartSlide presentation tool. Provide accurate, friendly, and concise responses about SmartSlide features and usage."},
                {"role": "user", "content": chatbot_prompt}
            ],
            "max_tokens": 1000,
            "temperature": 0.7
        }
        
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        
        result = response.json()
        bot_response = result["choices"][0]["message"]["content"]
        
        return jsonify({
            "response": bot_response,
            "status": "success"
        }), 200

    except requests.exceptions.HTTPError as http_err:
        current_app.logger.error(f"HTTP error during chatbot response: {http_err}")
        return jsonify({"error": "Failed to get response from AI service"}), 500
    except Exception as e:
        current_app.logger.error(f"Error in chatbot endpoint: {e}")
        return jsonify({"error": "Internal server error"}), 500
